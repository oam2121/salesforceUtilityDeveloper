import os
import pandas as pd
import plotly.express as px
import streamlit as st
from simple_salesforce import Salesforce
from reportlab.lib.units import inch
from st_aggrid import AgGrid
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Image
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
import tempfile

# Functions to fetch data and visualize

def get_salesforce_fields(sf, object_name):
    try:
        object_description = sf.__getattr__(object_name).describe()
        return [field['name'] for field in object_description['fields']]
    except Exception as e:
        st.error(f"Failed to fetch fields for {object_name}: {str(e)}")
        return []

def fetch_salesforce_data(sf, object_name, fields):
    soql_query = f"SELECT {', '.join(fields)} FROM {object_name}"
    results = sf.query_all(soql_query)
    return pd.json_normalize(results['records'])

def display_data_with_aggrid(df):
    AgGrid(df)

def choose_visualization(df, field):
    if df[field].dtype in ['float64', 'int64']:
        fig = px.histogram(df, x=field, title=f'Distribution of {field}', color_discrete_sequence=px.colors.qualitative.Plotly)
    elif df[field].dtype == 'object' and len(df[field].unique()) <= 20:
        fig = px.pie(df, names=field, title=f'Distribution of {field}', color_discrete_sequence=px.colors.qualitative.Plotly)
    else:
        fig = px.bar(df, x=field, title=f'Count of {field}', color_discrete_sequence=px.colors.qualitative.Plotly)
    
    st.plotly_chart(fig)
    return fig

# Updated save_plot_as_image function to maintain color scheme
def save_plot_as_image(fig):
    fig.update_layout(
        paper_bgcolor="white",  # Ensure paper background is white
        plot_bgcolor="white",   # Ensure plot background is white
        font=dict(color="black"),  # Ensure font color is black or another readable color
    )
    
    img_path = tempfile.NamedTemporaryFile(delete=False, suffix=".png").name
    fig.write_image(img_path, format='png', scale=2)  # High resolution
    return img_path

# Export to PDF with timestamp and visualizations
def export_to_pdf(df, figures):
    pdf_path = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf").name
    doc = SimpleDocTemplate(pdf_path, pagesize=letter)
    elements = []

    # Add timestamp
    timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    elements.append(Table([[f"Data Report generated on {timestamp}"]]))

    # Convert the dataframe to a list of lists for the table
    data = [df.columns.tolist()] + df.values.tolist()

    # Create a Table object
    table = Table(data)
    
    # Style the table
    table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 10),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
    ]))

    elements.append(table)

    # Add visualizations to the PDF
    for fig in figures:
        img_path = save_plot_as_image(fig)
        elements.append(Image(img_path, 6*inch, 4*inch))

    doc.build(elements)
    
    return pdf_path

def export_to_pptx(figures):
    pptx_path = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs = Presentation()
    for fig in figures:
        img_path = save_plot_as_image(fig)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(img_path, Inches(1), Inches(1), width=Inches(8))
    prs.save(pptx_path.name)
    return pptx_path.name

# Main smart visualize function
def smart_visualize(sf):
    st.subheader("Smart Visualize Salesforce Data")
    object_names = [obj["name"] for obj in sf.describe()["sobjects"]]
    selected_object = st.selectbox("Select Salesforce Object", object_names)

    fields = get_salesforce_fields(sf, selected_object)
    selected_fields = st.multiselect("Select fields to visualize", fields)

    show_data = st.button("Show Data")
    if show_data:
        df = fetch_salesforce_data(sf, selected_object, selected_fields)
        display_data_with_aggrid(df)
        figures = [choose_visualization(df, field) for field in selected_fields]

        # Save data and figures in session state to avoid recomputation
        st.session_state['data'] = df
        st.session_state['figures'] = figures

    # Ensure that we have both data and figures available in the session state
    if 'data' in st.session_state and 'figures' in st.session_state:
        # Generate PDF with data and figures
        if st.button("Export to PDF"):
            pdf_path = export_to_pdf(st.session_state['data'], st.session_state['figures'])
            with open(pdf_path, "rb") as file:
                btn = st.download_button(
                    label="Download PDF",
                    data=file,
                    file_name="report.pdf",
                    mime="application/pdf"
                )

        # Generate PPTX with figures
        if st.button("Export to PPTX"):
            pptx_path = export_to_pptx(st.session_state['figures'])
            with open(pptx_path, "rb") as file:
                btn = st.download_button(
                    label="Download PPTX",
                    data=file,
                    file_name="presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

if __name__ == "__main__":
    smart_visualize(sf)
